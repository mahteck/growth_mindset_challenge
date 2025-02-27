# Imports
import streamlit as st # type: ignore
import pandas as pd # type: ignore
import os
from docx import Document # type: ignore
from pptx import Presentation # type: ignore
import fitz  # PyMuPDF for PDF Handling # type: ignore
from fpdf import FPDF  # For PDF Generation # type: ignore
from io import BytesIO
import json

# ✅ Set Page Config at the start
st.set_page_config(page_title="💿 Data Sweeper", layout="wide")

# Sidebar - Settings
with st.sidebar:
    st.header("⚙️ Settings")
    theme_choice = st.radio("Choose Theme:", ["Light", "Dark"])
    st.info("ℹ️ Use buttons below for quick actions.")
    st.button("🔄 Refresh Page")

# ✅ Inject Custom CSS for Themes
def set_theme(theme_choice):
    if theme_choice == "Dark":
        dark_theme_css = """
        <style>
        body { background-color: #1e1e1e !important; color: white !important; }
        .stApp { background-color: #1e1e1e !important; color: white !important; }
        .stButton>button { background-color: #4CAF50 !important; color: white !important; }
        .stRadio>div { color: white !important; }
        h1, h2, h3, h4, h5, h6, p, label, .stAlert {
            color: white !important;
        }
        </style>
        """
        st.markdown(dark_theme_css, unsafe_allow_html=True)
    else:
        light_theme_css = """
        <style>
        body { background-color: white !important; color: black !important; }
        .stApp { background-color: white !important; color: black !important; }
        .stButton>button { background-color: #F8B501 !important; color: black !important; }
        .stRadio>div { color: black !important; }
        h1, h2, h3, h4, h5, h6, p, label, .stAlert {
            color: black !important;
        }

        /* ✅ Header Bar Light Theme */
        header { background-color: #F8B501 !important; }
        header * { color: black !important; }

        
        </style>
        """
        st.markdown(light_theme_css, unsafe_allow_html=True)

set_theme(theme_choice)  # ✅ Apply Theme


# App Title
st.markdown('<h1 style="color: #4CAF50;">💿 Data Sweeper</h1>', unsafe_allow_html=True)

st.write("Transform your files between (CSV, Excel, PowerPoint, PDF,  Word, Json) formats with built-in data cleaning and visualization!")

# File Uploader
# ✅ Additional Imports for Word and PowerPoint
from docx import Document # type: ignore
from pptx import Presentation # type: ignore

# ✅ Update File Uploader to Accept PDF, Word, json and PowerPoint Files
uploaded_files = st.file_uploader("Upload your files (CSV, Excel, Word, PowerPoint, PDF, Json):", 
                                  type=["csv", "xlsx", "docx", "pptx", "pdf", "Json"], 
                                  accept_multiple_files=True)

if uploaded_files:
    for file in uploaded_files:
        file_ext = os.path.splitext(file.name)[-1].lower()

        # ✅ Handle CSV and Excel
        if file_ext == ".csv":
            df = pd.read_csv(file)
        elif file_ext == ".xlsx":
            df = pd.read_excel(file)
        # ✅ Handle PDF Files
        elif file_ext == ".pdf":
            pdf_doc = fitz.open(stream=BytesIO(file.getvalue()))
            pdf_text = [page.get_text("text") for page in pdf_doc]
            df = pd.DataFrame({"Content": pdf_text})  

            st.success(f"✅ PDF File Loaded: {file.name}")
            st.subheader("📜 PDF Preview")
            st.text_area("PDF File Content", "\n\n".join(pdf_text), height=300)

        # ✅ Handle Word Files
        elif file_ext == ".docx":
            doc = Document(file)
            text = [para.text for para in doc.paragraphs if para.text.strip()]  
            df = pd.DataFrame({"Paragraphs": text})  

            st.success(f"✅ Word File Loaded: {file.name}")
            st.subheader("📄 Document Preview")
            st.text_area("Word File Content", "\n\n".join(text), height=300)
            
        # ✅ Handle PowerPoint Files
        elif file_ext == ".pptx":
            ppt = Presentation(file)
            slides_text = []
                
            for slide in ppt.slides:
                slide_text = "\n".join([shape.text for shape in slide.shapes if hasattr(shape, "text")])
                slides_text.append(slide_text)
                    
                df = pd.DataFrame({"Slides": slides_text})

            st.success(f"✅ PowerPoint File Loaded: {file.name}")
            st.subheader("📊 Presentation Preview")
            st.text_area("PowerPoint File Content", "\n\n--- Slide Break ---\n\n".join(slides_text), height=300)
        
        # ✅ Handle JSON Files
        elif file_ext == ".json":
            json_data = json.load(file)
            df = pd.json_normalize(json_data)

            st.success(f"✅ JSON File Loaded: {file_name}")
            st.subheader("🗂 JSON File Preview")
            st.json(json_data, expanded=True)

        else:
            st.error(f"❌ Unsupported file type: {file_ext}")
            continue  # Skip unsupported files
        
        # Display file details
        file_size_kb = len(file.getbuffer()) / 1024
        st.success(f"✅ File Loaded: {file.name} ({file_size_kb:.2f} KB)")

        # Show first 5 rows
        st.subheader("🔍 Data Preview")
        st.dataframe(df.head())

        # Data Cleaning Options
        st.subheader(f"🛠️ Data Cleaning Options")

        if st.checkbox(f"Clean Data for {file.name}"):
            col1, col2 = st.columns(2)

            # ✅ Fix: Use Session State to Persist Changes
            if f"remove_duplicates_{file.name}" not in st.session_state:
                st.session_state[f"remove_duplicates_{file.name}"] = False

            # with col1:
            #     if st.button(f"🚮 Remove Duplicates from {file.name}"):
            #         df.drop_duplicates(inplace=True)
            #         st.write("✅ Duplicates Removed Successfully!")

            with col1:
                    if st.button(f"🚮 Remove Duplicates from {file.name}"):
                        df.drop_duplicates(inplace=True)
                        st.session_state[f"remove_duplicates_{file.name}"] = True

            # ✅ Show Confirmation Message After Deduplication
            if st.session_state[f"remove_duplicates_{file.name}"]:
                st.success("✅ Duplicates Removed Successfully!")

            with col2:
                if st.button(f"🔧 Fill Missing Values for {file.name}"):
                    numeric_cols = df.select_dtypes(include=['number']).columns
                    df[numeric_cols] = df[numeric_cols].fillna(df[numeric_cols].mean())
                    st.success("✅ Missing Values Filled!")

        # Column Selection
        st.subheader("🎯 Choose Columns")
        columns = st.multiselect(f"Select Columns for {file.name}", df.columns, default=df.columns)    
        df = df[columns]

        # ✅ Data Visualization (Only for Numeric Columns)
        st.subheader("📊 Data Visualization")
        numeric_df = df.select_dtypes(include='number')
        if not numeric_df.empty:
            if st.checkbox(f"Show Graph for {file.name}"):
                st.bar_chart(numeric_df.iloc[:, :2])  # First 2 Numeric Columns
        else:
            st.warning(f"⚠️ No numeric data found in {file.name}. Graphs cannot be generated.")

        # Conversion Options
        st.subheader("🔁 Convert File Format")
        conversion_type = st.radio(
            f"Convert {file.name} to:", ["CSV", "Excel", "PDF", "PowerPoint", "Word", "JSON"], key=file.name
        )

        if st.button(f"📥 Convert {file.name}"):
            buffer = BytesIO()
    
            if conversion_type == "CSV":
                df.to_csv(buffer, index=False)
                file_name = file.name.replace(file_ext, ".csv")
                mime_type = "text/csv"
    
            elif conversion_type == "Excel":
                df.to_excel(buffer, index=False)
                file_name = file.name.replace(file_ext, ".xlsx")
                mime_type = "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet"
    
            elif conversion_type == "PDF":
                pdf = FPDF()
                pdf.set_auto_page_break(auto=True, margin=15)
                pdf.add_page()
                pdf.set_font("Arial", size=12)
    
                # Create a table header
                col_width = 40
                for col in df.columns:
                    pdf.cell(col_width, 10, col, border=1)
                pdf.ln()
    
                # Add table rows
                for _, row in df.iterrows():
                    for col in df.columns:
                        pdf.cell(col_width, 10, str(row[col]), border=1)
                    pdf.ln()

                pdf_buffer = BytesIO()
                pdf_output = pdf.output(dest="S").encode("latin1")  # Convert to bytes
                pdf_buffer.write(pdf_output)
                pdf_buffer.seek(0)
    
                file_name = file.name.replace(file_ext, ".pdf")
                buffer = pdf_buffer
                mime_type = "application/pdf"
    
            elif conversion_type == "Word":
                word_doc = Document()
                word_doc.add_heading("Data Export", level=1)
    
                table = word_doc.add_table(rows=1, cols=len(df.columns))
                hdr_cells = table.rows[0].cells
                for i, col in enumerate(df.columns):
                    hdr_cells[i].text = col
    
                for _, row in df.iterrows():
                    row_cells = table.add_row().cells
                    for i, col in enumerate(df.columns):
                        row_cells[i].text = str(row[col])
    
                word_buffer = BytesIO()
                word_doc.save(word_buffer)
                word_buffer.seek(0)
    
                file_name = file.name.replace(file_ext, ".docx")
                buffer = word_buffer
                mime_type = "application/vnd.openxmlformats-officedocument.wordprocessingml.document"
    
            elif conversion_type == "PowerPoint":
                ppt = Presentation()
                slide = ppt.slides.add_slide(ppt.slide_layouts[5])
                text_box = slide.shapes.add_textbox(10, 10, 700, 500)
                text_frame = text_box.text_frame
    
                text_frame.text = "\n".join(df.columns)
                for _, row in df.iterrows():
                    text_frame.text += "\n" + " | ".join(map(str, row.values))
    
                ppt_buffer = BytesIO()
                ppt.save(ppt_buffer)
                ppt_buffer.seek(0)
    
                file_name = file.name.replace(file_ext, ".pptx")
                buffer = ppt_buffer
                mime_type = "application/vnd.openxmlformats-officedocument.presentationml.presentation"
            
            elif conversion_type == "JSON":
                json_output = df.to_json(orient="records", indent=4)
                buffer.write(json_output.encode())
                buffer.seek(0)
                file_name = file.name.replace(file_ext, ".json")
                mime_type = "application/json"
            
            # Download Button
            st.download_button(
                label=f"📥 Download {file.name} as {conversion_type}",
                data=buffer,
                file_name=file_name,
                mime=mime_type
            )

st.success("🎉🥳 All files processed successfully!")  


# ✅ Footer with credit
st.markdown('<p style="text-align: center; font-size: 16px; color: #4CAF50;">Created by Shoaib Munir 🚀</p>', unsafe_allow_html=True)
